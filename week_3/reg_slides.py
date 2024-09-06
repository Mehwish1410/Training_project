from pptx import Presentation
from pptx.util import Pt

# Create a new PowerPoint presentation object
prs = Presentation()

# Function to add a slide with a title and bullet points
def add_bullet_slide(title, bullet_points):
    slide_layout = prs.slide_layouts[1]  # Use the "Title and Content" layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    tf = content_placeholder.text_frame
    tf.text = bullet_points[0]

    for point in bullet_points[1:]:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1  # Indent bullet point
        p.font.size = Pt(14)
    tf.word_wrap = True

# Content for the slides
slides_content = [
    ("Regression Analysis: An Overview", [
        "Regression analysis is a statistical method used to understand the relationship between a dependent variable (target) and one or more independent variables (predictors or features).",
        "The primary goal is to model and analyze the relationship to make predictions, identify trends, and infer causal relationships."
    ]),
    
    ("Basic Concepts", [
        "Dependent Variable (Target): The variable you are trying to predict or explain.",
        "Independent Variable (Predictors): The variables you are using to predict the dependent variable.",
        "Linear Regression: The simplest form of regression, where the relationship between the dependent and independent variables is modeled as a straight line.",
        "Regression Coefficient: Indicates the strength and direction of the relationship between an independent variable and the dependent variable.",
        "Intercept: The expected value of the dependent variable when all independent variables are zero.",
        "Residuals: The differences between observed values and the values predicted by the model.",
        "R-squared (R²): A statistical measure that represents the proportion of the variance for the dependent variable that's explained by the independent variables."
    ]),
    
    ("Objectives of Regression Analysis", [
        "Prediction: Predict the value of the dependent variable based on known values of independent variables.",
        "Estimation: Estimate the coefficients that quantify the relationship between the dependent and independent variables.",
        "Hypothesis Testing: Test hypotheses about the relationship between variables.",
        "Modeling Relationships: Understand the strength and nature of relationships between variables.",
        "Trend Analysis: Identify and analyze trends over time."
    ]),
    
    ("Types of Regression Analysis", [
        "Linear Regression",
        "  - Simple Linear Regression: Involves one independent variable.",
        "  - Multiple Linear Regression: Involves more than one independent variable.",
        "  - Example: Predicting house prices based on features like size, number of rooms, and location.",
        "Polynomial Regression",
        "  - Models the relationship as an nth degree polynomial.",
        "  - Example: Modeling the growth of a population where the relationship is nonlinear.",
        "Logistic Regression",
        "  - Used when the dependent variable is categorical (binary outcomes).",
        "  - Example: Classifying emails as spam or not spam.",
        "Ridge Regression",
        "  - A type of linear regression that includes a regularization term to prevent overfitting.",
        "  - Example: Predicting sales while controlling for multicollinearity among predictors.",
        "Lasso Regression",
        "  - Similar to Ridge Regression but can shrink some coefficients to zero, effectively selecting a simpler model.",
        "  - Example: Feature selection in high-dimensional datasets.",
        "Elastic Net Regression",
        "  - Combines both Ridge and Lasso regression penalties.",
        "  - Example: Used in complex models where feature selection and multicollinearity are both concerns.",
        "Stepwise Regression",
        "  - Automatically selects variables by adding or removing predictors based on specific criteria.",
        "  - Example: Building a predictive model with a large number of potential predictors.",
        "Quantile Regression",
        "  - Models the relationship between variables for different quantiles of the dependent variable.",
        "  - Example: Analyzing the impact of education on different percentiles of income distribution.",
        "Nonlinear Regression",
        "  - Models the relationship using nonlinear functions.",
        "  - Example: Growth curves in biology or pharmacokinetics."
    ]),
    
    ("Examples", [
        "Simple Linear Regression",
        "  - Objective: Predict the weight of a person based on their height.",
        "  - Model: Weight = b0 + b1 * Height",
        "  - Application: Understanding how height impacts weight.",
        "Multiple Linear Regression",
        "  - Objective: Predict house prices based on size, number of bedrooms, and age of the house.",
        "  - Model: Price = b0 + b1 * Size + b2 * Bedrooms + b3 * Age",
        "  - Application: Real estate pricing.",
        "Logistic Regression",
        "  - Objective: Determine the likelihood of a customer purchasing a product based on age, income, and browsing history.",
        "  - Model: Purchase Probability = 1 / (1 + e^-(b0 + b1 * Age + b2 * Income + b3 * BrowsingHistory))",
        "  - Application: Customer targeting in marketing.",
        "Polynomial Regression",
        "  - Objective: Model the relationship between hours studied and exam scores, which may not be linear.",
        "  - Model: Score = b0 + b1 * Hours + b2 * Hours^2",
        "  - Application: Educational performance analysis.",
        "Ridge Regression",
        "  - Objective: Predict stock prices with many correlated predictors.",
        "  - Model: Includes a penalty term to prevent overfitting.",
        "  - Application: Financial modeling."
    ]),
    
    ("Conclusion", [
        "Regression analysis is a fundamental tool in statistics and machine learning, enabling the modeling of relationships between variables.",
        "Understanding its basic concepts, objectives, and various types of regression allows for effective application across numerous fields, from finance and economics to biology and engineering."
    ])
]

# Add slides to the presentation
for title, bullet_points in slides_content:
    add_bullet_slide(title, bullet_points)

# Save the presentation
pptx_path = "Regression_Analysis.pptx"
prs.save(pptx_path)
